# -*- coding: utf-8 -*-
"""
Created on Thu Aug 26 17:50:19 2021

@author: sorgato
"""


def plot_dataf(machine, df, prs):
    import pandas as pd
    import matplotlib.pyplot as plt
    import numpy as np
    import pandas as pd
    import math
    from plot_inter import plot_interf
    from dose_corr import dose_corrf
    import seaborn as sns
    # from pptx import Presentation
    from pptx.util import Inches
    
    Doselimit=200
    Timelimit=100
    Countlimit=40
    matrix1= df

    # # global df, types_corr
    matrix = df[df['AMPLI'] == machine]
    # # matrix2 = df[df['AMPLI'] == machine[1]]
    # # matrix=matrix1.append(matrix2)
    # matrix = matrix1.sort_values(by='DATEINTERV').reset_index(drop=True)
    
    # # types = list(matrix['MOTIF'].unique())
    # types = list(matrix['MOTIF'])      
    # types_corr= tuple(types)
    # np.array(types)
    
    
    # dictionn= pd.read_excel ('dictionnaire_interventions_bloc+.xlsx', 'Feuil1')
    # com=list(dictionn.loc[:,'com'])
    # replace=list(dictionn.loc[:,'replace'])
    # categ=list(dictionn.loc[:,'catégorie'])
    
    # # from clean_types import clean_typesf
    # # types_corr=np.array(types)
    # # i=0
    # # for i in range(0, len(dictionn)):
    # #     com=list(dictionn.loc[i,'com'])
    # #     com=[''.join(map(str, com))]
    # #     replace= list(dictionn.loc[i,'replace'])
    # #     replace=[''.join(map(str, replace))]
    # #     types_corr =clean_typesf(com, replace, types_corr)
    # #     i=+i   
           

    # types_corr=[]
    # # i=0
    # for i in range(0, len(types)):
    #     typ=types[i]
    #     # s=14
    #     ii=0
    #     if not isinstance(typ, float):                 
    #         for s in range(0, len(com)):                
    #             if com[s] in typ and ii==0:
    #                 ii+=1
    #                 types_corr.append((typ, i, com[s], s, replace[s], categ[s]))
    #         if ii==0:
    #             types_corr.append((typ, i, 'lala', 0, 'lala', 'lala'))                
    #     else:
    #         types_corr.append((typ, i, 'nan', 0, 'nan', 'nan'))
                
    # types_corr=np.array(types_corr)
    
    # # ceux qui n'ont pas trouvé leur place dans dictionnaire
    # revoir = types_corr[np.array(np.where(types_corr[:,5] == 'lala')), :]
    
    # types_corr_uni=np.unique(types_corr[:,4])
    # # types_corr_uni.to_excel("dic_net.xlsx")
                      
    # # matrix['MOTIF_red']=types_corr[:,3]
    # # matrix['MOTIF_corr']=types_corr[:,4]
    # # matrix['CATEGORY']=types_corr[:,5]
    recap=pd.DataFrame()
    types_corr_uni=np.unique(matrix['MOTIF_corr'])
    recap['Intervention']=types_corr_uni
    recap['Median PDS (Gycm2)']=np.zeros(len(types_corr_uni))
    recap['75% percentile (Gycm2)']=np.zeros(len(types_corr_uni))
    recap['Mean time (s)']=np.zeros(len(types_corr_uni))
    recap['Count']=np.zeros(len(types_corr_uni))

    # i=0
  
    for i in range(0, len(types_corr_uni)):       
        t = types_corr_uni[i]
        mat = matrix.loc[matrix['MOTIF_corr'] == t]
        mat_dose = mat['DOSE Gycm2'] #µGym2        

        
        mat_time = mat['TEMPS (s)']
        # .fillna('0').astype(float)  # replace nan by 0
        # mat_time = [float('nan') if x == 0 else x for x in mat2]       
        
        
        # mat['DATEINTERV'] = mat['DATEINTERV'].astype('category')
        pd.plotting.register_matplotlib_converters()
        
        # try interactive plot: pour chaque année
        fig_int (mat, 2021)
        
       
        
        #calculate mean, 75% dose and time
        recap['Count'][i]=mat.shape[0]        
        recap['Median PDS (Gycm2)'][i]=round(np.nanmedian(mat_dose), 2)
        # recap['Mean PDS (Gycm2)'][i]=round(np.nanmean(mat_dose), 2)
        recap['75% percentile (Gycm2)'][i]=round(np.nanpercentile(mat_dose, 75), 2)
        recap['Mean time (s)'][i]=round(np.mean(mat_time), 2)

        
    # # plot histogram with mean PDS        # 
    # fig, ax = plt.subplots(1,1, figsize=(10,6))   
    # # ax = pd.cut(recap['Mean PDS (µGy m2)'][:-2], bins=len(types_corr_uni)-2, labels = recap['Intervention'][:-2]).value_counts(sort=False).plot.bar(y='Mean PDS (µGy m2)')
    # # ax = sns.countplot(x="Intervention", data=recap)   
    # ax = pd.cut(recap['Mean PDS (µGy m2)'][:-2], bins=len(types_corr_uni)-2, labels = recap['Intervention'][:-2]).plot.bar(y='Mean PDS (µGy m2)')
    # # plt.barrecap['Mean PDS (µGy m2)'][:-2], bins=len(types_corr_uni)-2, labels = recap['Intervention'][:-2]).value_counts(sort=False).plot.bar()
    # plt.ylabel('Mean PDS (µGy m²)')
    # plt.title(machine[1])
    # plt.tight_layout()
    
    # fig, ax = plt.subplots(1,1, figsize=(10,6))   
    # y=recap['Mean PDS (µGy m2)'][:-2]
    # x= recap['Intervention'][:-2]
    # width = 0.75 # the width of the bars 
    # ind = np.arange(len(types_corr_uni)-2)  # the x locations for the groups
    # ax.bar(ind, y, width, color="blue")
    # ax.set_xticks(ind)
    # ax.set_xticklabels(x, minor=False)
    # plt.xticks(rotation=90)
    # plt.title(machine)
    # plt.ylabel('Mean PDS (µGy m²)')     
    # plt.tight_layout()
    # fig.savefig('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures\\Moy_' + str(machine) +'.png') 

    recapCount=recap
    # Sort by higher to lower value  
    recapCount=recapCount.sort_values(by=['Count'], ascending=False)    
    recapCountE=recap.sort_values(by=['Count'], ascending=False)    
    FirstCount= recapCountE[recapCountE['Count'] > Countlimit]

    # plot bars with DF directly
    # recapDose=recap.drop('Mean time (s)', 1)
    recapDose=FirstCount[['Median PDS (Gycm2)', '75% percentile (Gycm2)', 'Count']]    
    # Sort by higher to lower value  
    recapDose=recapDose.sort_values(by=['Median PDS (Gycm2)'], ascending=False)
    recapDoseE=recap.sort_values(by=['Median PDS (Gycm2)'], ascending=False)
    FirstDose= recapDoseE[recapDoseE['Median PDS (Gycm2)'] > Doselimit]


    recapTime=FirstCount[['Mean time (s)', 'Count']]    
    # Sort by higher to lower value  
    recapTime=recapTime.sort_values(by=['Mean time (s)'], ascending=False)   
    recapTimeE=recap.sort_values(by=['Mean time (s)'], ascending=False)   
    FirstTime= recapTimeE[recapTimeE['Mean time (s)'] > Timelimit]
    
    


 # plot Dose       
    # fig=plt.figure()
    ax = recapDose[['Median PDS (Gycm2)', '75% percentile (Gycm2)']].plot(kind='bar', stacked=True, ylabel='PDS (Gycm²)', title=str(machine)+'>'+str(Countlimit), figsize=(10,6), ls='dashed', lw=3)#, facecolor="None")
    fig=ax.get_figure()
    ax.set_xticklabels(FirstCount.Intervention, rotation=90)     
    plt.tight_layout()# import pdb
# t=range(0, len(ax.patches), 2)
    for i in range(0, len(ax.patches)//2):
        # pdb.set_trace()
        # Find where everything is located
        rect= ax.patches[i]
        

        height = rect.get_height()
        width = rect.get_width()*2
        x = rect.get_x()
        y = rect.get_y()            
        # The height of the bar is the data value and can be used as the label
        label_text = f'{recapDose.Count.iloc[math.floor(i)]:.0f}'  # f'{height:.2f}' to format decimal values
        # ax.text(x, y, text)
        label_x = x + width / 4.
        label_y = y + height*1.5
    
        # plot only when height is greater than specified value
        # if height > 0:
        ax.text(label_x, label_y, label_text, ha='center', va='center', color='black', fontweight='bold')
        
 
    # add a little space at the top of the plot for the annotation
    ax.margins(y=0.1)    
    # move the legend out of the plot
    # ax.legend(title='Columns', bbox_to_anchor=(1, 1.02), loc='upper left')    
    plt.tight_layout()
    fig.savefig('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures\\Median_' + str(machine) +'.png') 


### make powerpoint ###


    # from pptx import Presentation
    # import os
   
    # prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    title = slide.shapes.title.text = str(machine)+'>'+str(Countlimit)
    # placeholder = slide.placeholders[1]
    pic = slide.shapes.add_picture('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures\\Median_' + str(machine) +'.png', Inches(1), Inches(1), width=Inches(10), height=Inches(5))
    # picture = placeholder.insert_picture('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures\\Moy_' + str(machine) +'.png', crop=False)
    # title = slide.shapes.title.text = str(machine)
    # _add_image(slide,1,"003.png")
# sub = slide.placeholders[2].text = "Python has the power"
# _add_image(slide,1,"003.png")
    # prs.save("NRLs.pptx")
    # os.startfile("NRLs.pptx")
 
    # return FirstDose FirstTime FirstCount