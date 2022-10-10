
###############################################################""""
#etude de dose par ampli: median, moyenne et 75%
# import pdb

dir="N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\spyder" 
os.chdir(dir)
# machine=amplis[1]
from plot_data import plot_dataf
import numpy as np
amplis=np.unique(df['AMPLI'].astype(str))
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title.text = 'NRLs'
for machine in amplis:
    # variab='DOSE'
    plot_dataf(machine, df, prs)
os.chdir('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures')
prs.save("NRLs.pptx")
os.startfile("NRLs.pptx")

from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title.text = 'NRLs'




# je suis restée là. Il faut creer une class First avec first dose, first time et firstcount donnés par plot_data pour chaque machine!


First
for machine in amplis:
    # machine='CIOS'
    plot_dataf(machine, df, prs)
    
    #trouver les intervention les plus irradiantes pour chaque ampli:
    FirstDose[] 
        
    
os.chdir('N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\figures')
prs.save("NRLs.pptx")
os.startfile("NRLs.pptx")


#############################################################################
#etude de dose par intervention et ampli: moyenne et 75%
dir="N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\spyder" 
os.chdir(dir)
intervs=np.unique(df['MOTIF_corr'].astype(str))
#trouver les 2 intervention les plus irradiantes pour chaque ampli:

    





#########################################################################  
#etude de gynécologues de la formation
dir="N:\\Themes\\Radioprotection GHM\\PYTHON_VSO\\NRLs\\spyder" 
os.chdir(dir)
from plot_data_chir import plot_dataf
# Ctrl+I in front of it, either on the Editor or the 
medecin=['Dr VERAN Camille', 'VERAN CAMILLE', 'Dr IOAN Bianca', 'IOAN BIANCA', 'Dr VIDAL Clementine', 'VIDAL CLEMENTINE', 'Dr FIZE Corinne', 'FIZE CORINNE', 'Dr GOUPIL Cecile', 'GOUPIL CECILE']
medecin=['Dr IOAN Bianca']
medecin=['Dr VIDAL Clementine']
medecin=['Dr FIZE Corinne']
medecin=['Dr GOUPIL Cecile', 'GOUPIL CECILE'] #rien!!
variab='DOSE'
df=dfTOT
plot_dataf(medecin,  dfTOT)